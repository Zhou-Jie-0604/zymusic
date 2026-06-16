#!/usr/bin/env python3
"""Generate ZYmusic technical roadmap PPT - polished version."""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ========== Colors ==========
C_DARK = RGBColor(0x1A, 0x3C, 0x6E)      # 深蓝
C_MED = RGBColor(0x2B, 0x57, 0x9A)       # 中蓝
C_LIGHT = RGBColor(0x3D, 0x7E, 0xC5)     # 浅蓝
C_ACCENT = RGBColor(0x5B, 0x9B, 0xD5)    # 亮蓝
C_DGREEN = RGBColor(0x1E, 0x6B, 0x4A)    # 深绿
C_MGREEN = RGBColor(0x2D, 0x8C, 0x5E)    # 中绿
C_LGREEN = RGBColor(0x4C, 0xAF, 0x50)    # 浅绿
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)    # 橙
C_LORANGE = RGBColor(0xF0, 0xA0, 0x50)   # 浅橙
C_PURPLE = RGBColor(0x7B, 0x3F, 0xA0)    # 紫
C_LPURPLE = RGBColor(0x9B, 0x6B, 0xC0)   # 浅紫
C_RED = RGBColor(0xC0, 0x39, 0x2B)       # 红
C_LRED = RGBColor(0xE7, 0x4C, 0x3C)      # 浅红
C_TEAL = RGBColor(0x00, 0x80, 0x80)      # 青
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x33, 0x33, 0x33)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_LGRAY = RGBColor(0xE8, 0xE8, 0xE8)
C_BG = RGBColor(0xF7, 0xF8, 0xFC)

bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = C_BG

def add_box(slide, left, top, width, height, fill, text='', fs=Pt(10), fc=C_WHITE, bold=False, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = fs
    tf.paragraphs[0].font.color.rgb = fc
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_text(slide, left, top, width, height, text, fs=Pt(9), fc=C_BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = fs
    tf.paragraphs[0].font.color.rgb = fc
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = align
    return tb

def add_multiline(slide, left, top, width, height, lines, fs=Pt(9), fc=C_WHITE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = fs; p.font.color.rgb = fc; p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return tb

def add_arrow_d(slide, x1, y1, x2, y2, color=C_GRAY, w=Pt(1.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = w
    return c

def add_v_arrow(slide, x, y, w, h, color=C_GRAY):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_layer_bg(slide, left, top, width, height, color, label):
    """Add a subtle background rectangle for a layer group"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Make it slightly transparent by using a lighter shade
    return shape

# ========== TITLE ==========
add_text(slide, Inches(0.5), Inches(0.12), Inches(15), Inches(0.55),
         'ZY音乐系统 — 技术路线图', Pt(24), C_DARK, True, PP_ALIGN.CENTER)

# Subtitle divider line
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(0.62), Inches(10), Pt(3))
div.fill.solid(); div.fill.fore_color.rgb = C_LIGHT; div.line.fill.background()

# ========== LAYER 1: 用户终端层 ==========
L1_Y = Inches(0.78)
add_text(slide, Inches(0.3), L1_Y, Inches(2.2), Inches(0.28), '▎用户终端层', Pt(11), C_DARK, True)

bx_w1 = Inches(3.8); bx_h1 = Inches(0.55); bx_gap = Inches(0.3)
bx_x_start = Inches(0.5)
bx_y1 = L1_Y + Inches(0.3)

boxes_L1 = [
    ('🖥 浏览器端 (Chrome / Edge / Firefox)', C_DARK),
    ('🖥 桌面客户端 (JavaFX WebView)', C_DGREEN),
    ('📱 移动端 PWA (Service Worker 离线)', C_PURPLE),
]
for i, (txt, clr) in enumerate(boxes_L1):
    add_box(slide, bx_x_start + i * (bx_w1 + bx_gap), bx_y1, bx_w1, bx_h1, clr, txt, Pt(13), C_WHITE, True)

# Down arrow L1 → L2
add_v_arrow(slide, Inches(7.6), bx_y1 + bx_h1 + Inches(0.02), Inches(0.45), Inches(0.3), C_MED)

# ========== LAYER 2: 表现层 + 控制层 ==========
L2_Y = bx_y1 + bx_h1 + Inches(0.42)
add_text(slide, Inches(0.3), L2_Y, Inches(2.5), Inches(0.28), '▎表现层 & 控制层', Pt(11), C_DARK, True)

l2_y = L2_Y + Inches(0.3); l2_h = Inches(0.70); l2_gap = Inches(0.2)
l2_boxes = [
    ('JSP 页面模板\n(index/play/upload/profile...)', C_LIGHT, Pt(10)),
    ('CSS3 样式 & 动画\n(响应式布局 / 粒子动画)', C_ACCENT, Pt(10)),
    ('JavaScript 交互脚本\n(播放控制 / 裁剪 / 异步)', C_TEAL, Pt(10)),
    ('Servlet 请求分发\n(路由映射 / 会话鉴权 / 转发)', C_MED, Pt(11)),
]
l2_w = Inches(2.85)
for i, (txt, clr, fs) in enumerate(l2_boxes):
    add_box(slide, Inches(0.5) + i * (l2_w + l2_gap), l2_y, l2_w, l2_h, clr, txt, fs, C_WHITE, True)

# Down arrow L2 → L3
add_v_arrow(slide, Inches(7.6), l2_y + l2_h + Inches(0.02), Inches(0.45), Inches(0.3), C_MED)

# ========== LAYER 3: 业务逻辑层 ==========
L3_Y = l2_y + l2_h + Inches(0.42)
add_text(slide, Inches(0.3), L3_Y, Inches(2.5), Inches(0.28), '▎业务逻辑层 (Service)', Pt(11), C_DARK, True)

l3_y = L3_Y + Inches(0.3); l3_h = Inches(0.70); l3_gap = Inches(0.12)
l3_data = [
    ('UserService', '用户认证\n会话管理', C_DARK),
    ('MusicService', '音乐上传\n点赞管理', C_MED),
    ('PlaylistService', '歌单CRUD\n排序管理', C_DGREEN),
    ('CommentService', '评论嵌套\n通知触发', C_MGREEN),
    ('PostService', '社区发帖\n内容管理', C_ORANGE),
    ('SearchService', '多维检索\n结果聚合', C_PURPLE),
]
l3_w = Inches(1.95)
for i, (title, desc, clr) in enumerate(l3_data):
    add_multiline(slide, Inches(0.5) + i * (l3_w + l3_gap), l3_y, l3_w, l3_h,
                  [(title, True), (desc, False)], Pt(9), C_WHITE)

# Down arrow L3 → L4
add_v_arrow(slide, Inches(7.6), l3_y + l3_h + Inches(0.02), Inches(0.45), Inches(0.3), C_MED)

# ========== LAYER 4: 数据访问层 + 存储层 ==========
L4_Y = l3_y + l3_h + Inches(0.42)
add_text(slide, Inches(0.3), L4_Y, Inches(3.0), Inches(0.28), '▎数据访问层 & 存储层', Pt(11), C_DARK, True)

l4_y = L4_Y + Inches(0.3); l4_h = Inches(0.65); l4_gap = Inches(0.2)
l4_data = [
    ('JDBC 数据访问\nPreparedStatement / 连接管理', C_MED),
    ('SQLite 关系数据库\nWAL日志 / 幂等迁移 / 并发控制', C_DGREEN),
    ('文件系统存储\nmusic / covers / avatars / backgrounds', C_ORANGE),
    ('嵌入式 Jetty 服务器\n端口自适应 / HTTP Range / 多部件上传', C_PURPLE),
]
l4_w = Inches(2.95)
for i, (txt, clr) in enumerate(l4_data):
    add_box(slide, Inches(0.5) + i * (l4_w + l4_gap), l4_y, l4_w, l4_h, clr, txt, Pt(11), C_WHITE, True)

# ========== RIGHT SIDEBAR: 后端核心算法 ==========
ALGO_R_X = Inches(13.1); ALGO_W = Inches(2.7)
ALGO_START = Inches(0.85)

add_text(slide, ALGO_R_X, Inches(0.78), ALGO_W, Inches(0.28), '▎后端核心算法', Pt(10), C_RED, True, PP_ALIGN.CENTER)

backend_algos = [
    ('邻接表模型', '评论嵌套树存储与查询', C_RED),
    ('LIKE 通配符模式匹配', '四维度模糊音乐检索', C_LRED),
    ('Toggle 原子状态翻转', '点赞检查-切换-计数', C_DARK),
    ('HTTP Range 分段传输', '字节级流式定位播放', C_MED),
    ('元数据列检测迁移', '幂等增量数据库演进', C_PURPLE),
]

for i, (title, desc, clr) in enumerate(backend_algos):
    ay = ALGO_START + i * Inches(0.7)
    add_box(slide, ALGO_R_X, ay, ALGO_W, Inches(0.55), clr, '', Pt(8), C_WHITE, True)
    add_multiline(slide, ALGO_R_X, ay + Inches(0.02), ALGO_W, Inches(0.5),
                  [(title, True), (desc, False)], Pt(8), C_WHITE)

# ========== LEFT SIDEBAR: 前端核心算法 ==========
FE_ALGO_X = Inches(0.15); FE_ALGO_W = Inches(2.5)
FE_START = Inches(5.65)

add_text(slide, FE_ALGO_X, FE_START - Inches(0.25), FE_ALGO_W, Inches(0.25),
         '▎前端核心算法', Pt(10), C_DGREEN, True, PP_ALIGN.CENTER)

fe_algos = [
    ('IIFE 单例模式', '全局播放器跨页面共享', C_DGREEN),
    ('有限状态机', '循环/顺序/随机模式控制', C_MGREEN),
    ('Canvas 仿射变换', '圆形头像裁剪交互管线', C_LGREEN),
    ('Cache-First 策略', 'PWA 离线缓存与请求拦截', C_TEAL),
    ('递归树遍历', '评论扁平-嵌套结构转换', C_PURPLE),
]

for i, (title, desc, clr) in enumerate(fe_algos):
    ay = FE_START + i * Inches(0.58)
    add_box(slide, FE_ALGO_X, ay, FE_ALGO_W, Inches(0.47), clr, '', Pt(7), C_WHITE, True)
    add_multiline(slide, FE_ALGO_X, ay + Inches(0.02), FE_ALGO_W, Inches(0.43),
                  [(title, True), (desc, False)], Pt(7), C_WHITE)

# ========== DATA FLOW: big vertical arrow line ==========
# Left side vertical line connecting backend algos to main flow
for i in range(len(backend_algos)):
    ay = ALGO_START + i * Inches(0.7) + Inches(0.27)
    add_arrow_d(slide, ALGO_R_X, ay, Inches(12.8), ay, C_LGRAY, Pt(0.8))

for i in range(len(fe_algos)):
    ay = FE_START + i * Inches(0.58) + Inches(0.23)
    add_arrow_d(slide, FE_ALGO_X + FE_ALGO_W, ay, Inches(0.5), ay, C_LGRAY, Pt(0.8))

# ========== BOTTOM: 技术栈汇总 ==========
TECH_Y = Inches(7.75)
add_text(slide, Inches(0.3), TECH_Y, Inches(3), Inches(0.25), '▎关键技术栈', Pt(10), C_DARK, True)

tech_items = [
    ('Java 17', C_DARK), ('Servlet 4.0', C_MED), ('JSP 2.3', C_LIGHT),
    ('Jetty 9.4', C_ACCENT), ('JDBC', C_DGREEN), ('SQLite 3', C_MGREEN),
    ('HTML5', C_ORANGE), ('CSS3', C_LORANGE), ('JavaScript ES6', C_PURPLE),
    ('Canvas API', C_LPURPLE), ('Fetch API', C_TEAL), ('Service Worker', C_DARK),
    ('JavaFX', C_MED), ('Maven 3', C_DGREEN), ('Gson', C_LIGHT),
]
tech_w = Inches(1.55); tech_h = Inches(0.38); tech_gap = Inches(0.08)
items_per_row = 7
for i, (name, clr) in enumerate(tech_items):
    col = i % items_per_row
    row = i // items_per_row
    x = Inches(0.5) + col * (tech_w + tech_gap)
    y = TECH_Y + Inches(0.28) + row * (tech_h + Inches(0.06))
    add_box(slide, x, y, tech_w, tech_h, clr, name, Pt(9), C_WHITE, True, radius=False)

# ========== FOOTER ==========
add_text(slide, Inches(0.5), Inches(8.7), Inches(15), Inches(0.25),
         'ZY音乐系统  |  Java Web 全栈架构  |  浏览器 + 桌面 + 移动端 三端统一', Pt(8), C_GRAY, False, PP_ALIGN.CENTER)

# ========== SAVE ==========
output = 'F:/ZYmusic(2)/ZY音乐技术路线图.pptx'
prs.save(output)
print(f'Saved: {output}')
